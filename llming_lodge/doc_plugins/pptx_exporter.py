"""Server-side PPTX generation using python-pptx with a real template file.

Dual path:
- **Template-native**: slide has ``placeholders`` dict + template_config has ``layouts`` →
  uses layout_index and ph_idx from the config to fill real template placeholders.
- **Legacy (abstract)**: slide has ``elements`` list → old heuristic (title/content/end detection).
"""

import base64
import io
import logging
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)

# ── Legacy layout indices (fallback when no template_config) ──────
LAYOUT_TITLE = 1
LAYOUT_TEXT = 6
LAYOUT_END = 15

OBJ_LEFT = Inches(0.39)
OBJ_TOP = Inches(1.54)
OBJ_WIDTH = Inches(13.19)
OBJ_HEIGHT = Inches(5.79)


# ── Shared helpers ────────────────────────────────────────────────

def _elem_text(elem: dict) -> str:
    """Get text content from an element (handles both 'content' and 'text' field names)."""
    return elem.get("content") or elem.get("text") or ""


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#RRGGBB' to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_text(placeholder, text: str, font_size: int = 18,
              bold: bool = False, color: str = "", align: str = "",
              font_name: str = "Arial"):
    """Set text content on a placeholder with formatting."""
    placeholder.text = ""
    tf = placeholder.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0] if p.runs else p.add_run()
    if not p.runs:
        run.text = text
    run.font.size = Pt(font_size)
    run.font.name = font_name
    run.font.bold = bold
    if color:
        run.font.color.rgb = _hex_to_rgb(color)
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT


def _add_paragraphs(text_frame, items: list[str], font_size: int = 16,
                    color: str = "#333333", font_name: str = "Arial",
                    bullet: bool = True):
    """Add multiple paragraphs (e.g. bullet items) to a text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = str(item)
        p.font.size = Pt(font_size)
        p.font.name = font_name
        if color:
            p.font.color.rgb = _hex_to_rgb(color)
        if bullet:
            p.level = 0


def _add_image_from_base64(slide, img_data: str, left, top, width, height):
    """Add a picture from a base64 data URI or raw base64 string."""
    if img_data.startswith("data:"):
        img_data = img_data.split(",", 1)[1]
    img_bytes = base64.b64decode(img_data)
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, left, top, width, height)


def _add_table_shape(slide, headers: list, rows: list,
                     left, top, width, height,
                     accent_color: str = "#1D459F", font_name: str = "Arial"):
    """Add a table shape at the given geometry."""
    if not headers and not rows:
        return
    cols = len(headers) if headers else (len(rows[0]) if rows else 0)
    total_rows = (1 if headers else 0) + len(rows)
    tbl = slide.shapes.add_table(total_rows, cols, left, top, width, height).table
    if headers:
        for c, h_text in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = str(h_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.name = font_name
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex_to_rgb(accent_color)
    row_start = 1 if headers else 0
    for r, row_data in enumerate(rows):
        if not isinstance(row_data, list):
            continue
        for c, cell_val in enumerate(row_data):
            if c >= cols:
                break
            cell = tbl.cell(row_start + r, c)
            cell.text = str(cell_val) if cell_val is not None else ""
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = font_name


def _remove_existing_slides(prs):
    """Remove all slides from a presentation (clean slate for generation)."""
    _R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    while len(prs.slides) > 0:
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.get(_R_NS + "id")
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)


# ── Template-native path ─────────────────────────────────────────

def _find_layout_def(template_config: dict, layout_name: str) -> dict | None:
    """Find a layout definition by name in the template config."""
    for layout in template_config.get("layouts", []):
        if layout.get("name") == layout_name:
            return layout
    return None


def _fill_placeholder_value(slide, ph_idx: int, value, ph_def: dict,
                            chart_images: dict, accent_color: str,
                            text_color: str, subtitle_color: str,
                            font_name: str,
                            is_title_layout: bool = False,
                            is_end_layout: bool = False):
    """Fill a single placeholder with its value (string or rich object).

    Uses font sizes and colors matching the JS preview — not the template's
    native sizes (which can be wildly different, e.g. 68pt for titles).
    """
    # Find the placeholder on the slide by ph_idx
    ph = None
    for p in slide.placeholders:
        if p.placeholder_format.idx == ph_idx:
            ph = p
            break

    ph_name = ph_def.get("name", "")

    if isinstance(value, str):
        if ph:
            # Font sizes matching JS _layoutTemplateSlide()
            is_title = ph_name == "title"
            is_body = ph_name in ("body", "subtitle")
            if is_title:
                font_size = 36 if is_title_layout else 24
                color = "#FFFFFF" if is_title_layout else accent_color
                bold = True
            elif is_body:
                font_size = 14
                color = subtitle_color
                bold = False
            else:
                font_size = 16
                color = text_color
                bold = False
            align = "center" if (is_title_layout or is_end_layout) else ""
            _set_text(ph, value, font_size=font_size, bold=bold,
                      color=color, align=align, font_name=font_name)
        return

    if not isinstance(value, dict):
        return

    vtype = value.get("type", "")
    # Geometry from placeholder def (inches → Emu)
    left = Inches(ph_def.get("x", 0))
    top = Inches(ph_def.get("y", 0))
    width = Inches(ph_def.get("w", 10))
    height = Inches(ph_def.get("h", 5))

    if vtype == "list":
        items = value.get("items", [])
        if ph and hasattr(ph, "text_frame"):
            _add_paragraphs(ph.text_frame, items, font_size=16,
                            color=text_color, font_name=font_name)
        return

    if vtype == "table":
        # Tables can't go into a text placeholder — add as a shape
        # Delete the placeholder to make room if it exists
        if ph:
            sp = ph._element
            sp.getparent().remove(sp)
        _add_table_shape(
            slide, value.get("headers", []), value.get("rows", []),
            left, top, width, height,
            accent_color=accent_color, font_name=font_name,
        )
        return

    if vtype == "chart":
        chart_id = value.get("_chartImageId", "")
        img_data = chart_images.get(chart_id, "")
        if not img_data:
            return
        # Remove placeholder to make room
        if ph:
            sp = ph._element
            sp.getparent().remove(sp)
        try:
            _add_image_from_base64(slide, img_data, left, top, width, height)
        except Exception as e:
            logger.warning(f"[PPTX_EXPORT] Failed to add chart image: {e}")
        return

    if vtype == "image":
        src = value.get("src", "")
        if not src:
            return
        if ph:
            sp = ph._element
            sp.getparent().remove(sp)
        try:
            _add_image_from_base64(slide, src, left, top, width, height)
        except Exception as e:
            logger.warning(f"[PPTX_EXPORT] Failed to add image: {e}")
        return

    # Fallback: treat as text
    text = value.get("content") or value.get("text") or str(value)
    if ph:
        _set_text(ph, text, font_size=16, font_name=font_name, color=text_color)


def _add_template_native_slide(prs, slide_spec: dict, template_config: dict,
                                chart_images: dict, accent_color: str,
                                text_color: str, subtitle_color: str,
                                font_name: str):
    """Add a slide using the template-native format (layout name + placeholders dict)."""
    layout_name = slide_spec.get("layout", "text")
    layout_def = _find_layout_def(template_config, layout_name)
    if not layout_def:
        logger.warning(f"[PPTX_EXPORT] Layout '{layout_name}' not found in template config")
        return

    layout_index = layout_def.get("layoutIndex", 0)
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)

    placeholders = slide_spec.get("placeholders", {})
    ph_defs = {pd["name"]: pd for pd in layout_def.get("placeholders", [])}

    is_title_layout = layout_def.get("isTitle", False)
    is_end_layout = layout_def.get("isEnd", False)

    for ph_name, value in placeholders.items():
        ph_def = ph_defs.get(ph_name)
        if not ph_def:
            logger.debug(f"[PPTX_EXPORT] Unknown placeholder '{ph_name}' in layout '{layout_name}'")
            continue

        _fill_placeholder_value(
            slide, ph_def["phIdx"], value, ph_def,
            chart_images, accent_color, text_color, subtitle_color,
            font_name,
            is_title_layout=is_title_layout,
            is_end_layout=is_end_layout,
        )


# ── Legacy (abstract) path ───────────────────────────────────────

def _is_title_slide(slide_spec: dict, idx: int) -> bool:
    if slide_spec.get("layout") == "title":
        return True
    if idx != 0:
        return False
    elems = slide_spec.get("elements", [])
    if len(elems) == 0:
        return True
    if len(elems) == 1 and elems[0].get("type") in ("text", "subtitle"):
        return True
    return False


def _is_end_slide(slide_spec: dict, idx: int, total: int) -> bool:
    if slide_spec.get("layout") == "end":
        return True
    if idx != total - 1 or idx == 0:
        return False
    elems = slide_spec.get("elements", [])
    if len(elems) == 0:
        return True
    if len(elems) == 1 and elems[0].get("type") in ("text", "subtitle"):
        return True
    return False


def _fill_content_elements(slide, elements: list[dict], chart_images: dict,
                           accent_color: str = "#1D459F",
                           text_color: str = "#333333",
                           font_name: str = "Arial"):
    """Fill content elements into the OBJECT placeholder area of a content slide."""
    obj_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 17:
            obj_ph = ph
            break

    y_offset = OBJ_TOP
    text_elements = []
    shape_elements = []

    for elem in elements:
        etype = elem.get("type", "")
        if etype in ("text", "heading", "subtitle", "list"):
            text_elements.append(elem)
        elif etype in ("table", "chart", "image"):
            shape_elements.append(elem)

    if obj_ph and text_elements:
        tf = obj_ph.text_frame
        tf.word_wrap = True
        first = True
        for elem in text_elements:
            etype = elem.get("type", "")
            if etype == "heading":
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = _elem_text(elem)
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.name = font_name
                p.font.color.rgb = _hex_to_rgb(accent_color)
                p.space_after = Pt(6)
            elif etype == "list":
                items = elem.get("items", [])
                if not items and _elem_text(elem):
                    items = _elem_text(elem).split("\n")
                for item in items:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = item
                    p.font.size = Pt(16)
                    p.font.name = font_name
                    p.font.color.rgb = _hex_to_rgb(text_color)
                    p.level = 0
                    p.space_after = Pt(4)
            else:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = _elem_text(elem)
                p.font.size = Pt(16)
                p.font.name = font_name
                p.font.color.rgb = _hex_to_rgb(text_color)
                p.space_after = Pt(4)

        text_line_count = sum(
            1 + (len(e.get("items", [])) if e.get("type") == "list" else 0)
            for e in text_elements
        )
        y_offset = OBJ_TOP + Emu(int(text_line_count * Pt(24)))
    elif not text_elements:
        y_offset = OBJ_TOP

    for elem in shape_elements:
        etype = elem.get("type", "")
        remaining_h = (OBJ_TOP + OBJ_HEIGHT) - y_offset
        if remaining_h < Inches(0.5):
            break

        if etype == "table":
            headers = elem.get("headers", [])
            rows = elem.get("rows", [])
            if not headers and not rows:
                continue
            cols = len(headers) if headers else (len(rows[0]) if rows else 0)
            total_rows = (1 if headers else 0) + len(rows)
            table_h = min(remaining_h, Inches(total_rows * 0.35 + 0.3))
            _add_table_shape(
                slide, headers, rows,
                OBJ_LEFT, y_offset, OBJ_WIDTH, table_h,
                accent_color=accent_color, font_name=font_name,
            )
            y_offset += table_h + Inches(0.2)

        elif etype == "chart":
            chart_id = elem.get("_chartImageId", "")
            img_data = chart_images.get(chart_id, "")
            if not img_data:
                continue
            try:
                chart_h = min(remaining_h, Inches(4.5))
                _add_image_from_base64(slide, img_data, OBJ_LEFT, y_offset, OBJ_WIDTH, chart_h)
                y_offset += chart_h + Inches(0.2)
            except Exception as e:
                logger.warning(f"[PPTX_EXPORT] Failed to add chart image: {e}")

        elif etype == "image":
            src = elem.get("src", "")
            if not src:
                continue
            try:
                img_h = min(remaining_h, Inches(4.0))
                _add_image_from_base64(slide, src, OBJ_LEFT, y_offset, OBJ_WIDTH, img_h)
                y_offset += img_h + Inches(0.2)
            except Exception as e:
                logger.warning(f"[PPTX_EXPORT] Failed to add image: {e}")


def _add_legacy_slide(prs, slide_spec: dict, idx: int, total: int,
                      chart_images: dict, accent: str, text_color: str,
                      subtitle_color: str, font_name: str):
    """Add a slide using the legacy abstract format (elements list)."""
    is_title = _is_title_slide(slide_spec, idx)
    is_end = _is_end_slide(slide_spec, idx, total)

    if is_title:
        layout = prs.slide_layouts[LAYOUT_TITLE]
        slide = prs.slides.add_slide(layout)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                _set_text(ph, slide_spec.get("title", ""),
                          font_size=36, bold=True, color="#FFFFFF",
                          align="center", font_name=font_name)
            elif ph.placeholder_format.idx == 1:
                sub_elem = next(
                    (e for e in slide_spec.get("elements", [])
                     if e.get("type") in ("subtitle", "text")),
                    None,
                )
                sub_text = _elem_text(sub_elem) if sub_elem else ""
                if sub_text:
                    _set_text(ph, sub_text, font_size=20,
                              color="#FFFFFF", align="center",
                              font_name=font_name)
                else:
                    ph.text = ""

    elif is_end:
        layout = prs.slide_layouts[LAYOUT_END]
        slide = prs.slides.add_slide(layout)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                _set_text(ph, slide_spec.get("title", ""),
                          font_size=32, bold=True, color=accent,
                          align="center", font_name=font_name)
            elif ph.placeholder_format.idx == 10:
                sub_elem = next(
                    (e for e in slide_spec.get("elements", [])
                     if e.get("type") in ("subtitle", "text")),
                    None,
                )
                if sub_elem:
                    _set_text(ph, _elem_text(sub_elem),
                              font_size=18, color=subtitle_color,
                              align="center", font_name=font_name)
                else:
                    ph.text = ""

    else:
        layout = prs.slide_layouts[LAYOUT_TEXT]
        slide = prs.slides.add_slide(layout)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                _set_text(ph, slide_spec.get("title", ""),
                          font_size=24, bold=True, color=accent,
                          font_name=font_name)
            elif ph.placeholder_format.idx == 16:
                sub_elem = next(
                    (e for e in slide_spec.get("elements", [])
                     if e.get("type") == "subtitle"),
                    None,
                )
                if sub_elem:
                    _set_text(ph, _elem_text(sub_elem),
                              font_size=14, color=subtitle_color,
                              font_name=font_name)
                else:
                    ph.text = ""

        content_elems = [
            e for e in slide_spec.get("elements", [])
            if e.get("type") != "subtitle"
        ]
        if content_elems:
            _fill_content_elements(
                slide, content_elems, chart_images,
                accent_color=accent,
                text_color=text_color,
                font_name=font_name,
            )


# ── Public API ────────────────────────────────────────────────────

def export_pptx(spec: dict, template_path: str,
                chart_images: dict | None = None,
                template_config: dict | None = None) -> bytes:
    """Generate PPTX from spec using the real template file.

    Args:
        spec: Presentation spec from the frontend.
        template_path: Path to the .pptx template file.
        chart_images: Dict mapping chart element IDs to base64 PNG data.
        template_config: Serialized PresentationTemplate with layouts (camelCase).
            When provided and slides use ``placeholders``, the template-native
            path is used.  Otherwise falls back to the legacy path.

    Returns:
        PPTX file content as bytes.
    """
    chart_images = chart_images or {}
    template_config = template_config or {}
    has_layouts = bool(template_config.get("layouts"))

    prs = Presentation(template_path)
    slides_spec = spec.get("slides", [])
    total = len(slides_spec)
    accent = template_config.get("accentColor") or spec.get("accentColor", "#1D459F")
    text_color = template_config.get("textColor") or spec.get("textColor", "#333333")
    subtitle_color = template_config.get("subtitleColor") or spec.get("subtitleColor", "#5695B4")
    font_name = template_config.get("headingFont") or spec.get("headingFont", "Arial")

    _remove_existing_slides(prs)

    for i, slide_spec in enumerate(slides_spec):
        if has_layouts and slide_spec.get("placeholders"):
            _add_template_native_slide(
                prs, slide_spec, template_config,
                chart_images, accent, text_color, subtitle_color, font_name,
            )
        else:
            _add_legacy_slide(
                prs, slide_spec, i, total,
                chart_images, accent, text_color, subtitle_color, font_name,
            )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
